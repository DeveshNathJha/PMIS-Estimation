"""
PMIS Alpha Prototype - Premium Presentation Generator v2
==========================================================
Dual-audience layout:
    LEFT (60 %) - detailed technical bullet content
    RIGHT (38 %) - "What This Means"plain-English panel
    BOTTOM strip - Key Insight / Takeaway callout

Design System:
    Navy #0D1B2A · Electric Blue #1E90FF · Cyan #00C8FF
    White #F0F4F8 · Muted #A0AFBE · Amber #F5A623
    Green #27AE60 · Red #E74C3C · Jet #07090D

Run:
    pip3 install python-pptx --break-system-packages
    python3 generate_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Dimensions (16:9 widescreen) 
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Palette 
C_JET = RGBColor(0x07, 0x09, 0x0D)
C_NAVY = RGBColor(0x0D, 0x1B, 0x2A)
C_PANEL = RGBColor(0x12, 0x22, 0x35)
C_PANEL2 = RGBColor(0x0A, 0x17, 0x26)
C_ACCENT = RGBColor(0x1E, 0x90, 0xFF)
C_CYAN = RGBColor(0x00, 0xC8, 0xFF)
C_WHITE = RGBColor(0xF0, 0xF4, 0xF8)
C_MUTED = RGBColor(0xA0, 0xAF, 0xBE)
C_AMBER = RGBColor(0xF5, 0xA6, 0x23)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)

# Low-level helpers 

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
    size=12, bold=False, italic=False,
    color=None, align=PP_ALIGN.LEFT,
    font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return tb


def rule(slide, l, t, w, color=None):
    r = slide.shapes.add_shape(1, l, t, w, Inches(0.018))
    r.fill.solid()
    r.fill.fore_color.rgb = color or C_ACCENT
    r.line.fill.background()
    return r


def badge(slide, text, l, t, w=Inches(2.2), h=Inches(0.28),
    bg=None, fg=None, size=8):
    bg = bg or C_ACCENT
    fg = fg or C_JET
    s = rect(slide, l, t, w, h, bg)
    tf = s.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(size)
    run.font.color.rgb = fg
    run.font.name = "Calibri"
    return s


def brand_bar(slide, num, total=14):
    rect(slide, 0, SLIDE_H - Inches(0.27), SLIDE_W, Inches(0.27), C_NAVY)
    txt(slide,
    "PMIS Alpha Prototype | For Internal Review Only",
    Inches(0.3), SLIDE_H - Inches(0.26),
    SLIDE_W - Inches(2.5), Inches(0.26),
    size=7, color=C_MUTED, align=PP_ALIGN.LEFT)
    txt(slide,
    f"Slide {num:02d} / {total}",
    SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.26),
    Inches(1.1), Inches(0.26),
    size=7, color=C_MUTED, align=PP_ALIGN.RIGHT)


def left_bar(slide, color=None):
    rect(slide, 0, 0, Inches(0.18), SLIDE_H, color or C_ACCENT)


def slide_header(slide, num, icon, title, subtitle, accent=None):
    ac = accent or C_ACCENT
    rect(slide, 0, 0, SLIDE_W, Inches(1.5), C_NAVY)
    left_bar(slide, ac)
    badge(slide, f"{icon} SLIDE {num:02d} ",
    Inches(0.35), Inches(0.1), Inches(1.7), Inches(0.27),
    bg=ac, fg=C_JET, size=8)
    txt(slide, title,
    Inches(0.38), Inches(0.38), Inches(12.0), Inches(0.72),
    size=27, bold=True, color=C_WHITE)
    txt(slide, subtitle,
    Inches(0.38), Inches(1.1), Inches(12.0), Inches(0.36),
    size=12, italic=True, color=C_CYAN)
    rule(slide, Inches(0.38), Inches(1.48), SLIDE_W - Inches(0.5), color=ac)

# Multi-paragraph rich text box 
def rich_bullets(slide, items, l, t, w, h):
    """
    items: list of (label, body_text) tuples.
    label can be '' for a plain bullet.
    Each bullet auto-wraps. Spacing is controlled via paragraph spacing.
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for label, body in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)
        if label:
            r1 = p.add_run()
            r1.text = f"  {label}  "
            r1.font.bold = True
            r1.font.size = Pt(11.5)
            r1.font.color.rgb = C_ACCENT
            r1.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = body if label else f"  {body}"
        r2.font.bold = False
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = C_WHITE
        r2.font.name = "Calibri"
    return tb


def plain_english_panel(slide, lines, l, t, w, h, accent=None):
    """
    Right-side 'In Simple Terms' panel for non-technical audience.
    lines: list of strings (plain language sentences)
    """
    ac = accent or C_CYAN
    # panel background
    rect(slide, l, t, w, h, C_PANEL2)
    # top accent stripe
    rect(slide, l, t, w, Inches(0.07), ac)
    # label
    txt(slide, "IN SIMPLE TERMS",
        l + Inches(0.1), t + Inches(0.1),
        w - Inches(0.2), Inches(0.28),
        size=8, bold=True, color=ac)
    rule(slide, l + Inches(0.1), t + Inches(0.38),
         w - Inches(0.2), color=ac)
    # content
    tb = slide.shapes.add_textbox(
        l + Inches(0.12), t + Inches(0.48),
        w - Inches(0.22), h - Inches(0.58))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(7)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(11)
        run.font.color.rgb = C_WHITE
        run.font.name = "Calibri"
        run.font.italic = False
    return tb



def key_insight(slide, text, l, t, w, accent=None):
    """Bottom-strip key insight / takeaway for all audiences."""
    ac = accent or C_ACCENT
    h = Inches(0.6)
    rect(slide, l, t, w, h, RGBColor(0x08, 0x1A, 0x2E))
    rect(slide, l, t, Inches(0.07), h, ac)
    tb = slide.shapes.add_textbox(
    l + Inches(0.18), t + Inches(0.06),
    w - Inches(0.25), h - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "KEY INSIGHT: "
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.color.rgb = ac
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(10)
    r2.font.color.rgb = C_WHITE
    r2.font.name = "Calibri"
    return tb


def constraint_box(slide, text, l, t, w):
    """Amber constraint warning box."""
    h = Inches(0.52)
    rect(slide, l, t, w, h, RGBColor(0x2A, 0x12, 0x00))
    rect(slide, l, t, Inches(0.07), h, C_AMBER)
    tb = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.07),
    w - Inches(0.25), h - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "NOTE: "
    r1.font.bold = True
    r1.font.size = Pt(9.5)
    r1.font.color.rgb = C_AMBER
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(9.5)
    r2.font.color.rgb = C_WHITE
    r2.font.name = "Calibri"
    return tb


# 
# DUAL-AUDIENCE CONTENT SLIDE BUILDER
# 

LEFT_W = Inches(7.6)
RIGHT_W = Inches(4.8)
LEFT_L = Inches(0.3)
RIGHT_L = Inches(8.2)
CONTENT_TOP = Inches(1.6)

def dual_slide(prs, num, icon, title, subtitle,
    tech_bullets, # list of (label, detailed_body)
    plain_lines, # list of plain-English sentences
    insight, # key insight string
    constraint=None, # optional constraint string
    accent=None):
    """
    Builds a dual-audience slide:
    Left (60%) : technical expanded bullets
    Right (38%) : plain-English panel
    Bottom : key insight strip (+ optional constraint above it)
    """
    ac = accent or C_ACCENT
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_JET)
    slide_header(slide, num, icon, title, subtitle, ac)

    # work out vertical space
    bottom_reserved = Inches(0.28)  # brand bar
    insight_h = Inches(0.6)
    constraint_h = Inches(0.52) + Inches(0.08) if constraint else 0
    content_bot = SLIDE_H - bottom_reserved - insight_h - constraint_h - Inches(0.1)
    content_h = content_bot - CONTENT_TOP

    # Left technical panel
    rect(slide, LEFT_L, CONTENT_TOP, LEFT_W, content_h, C_PANEL)
    rich_bullets(slide, tech_bullets,
                 LEFT_L + Inches(0.18), CONTENT_TOP + Inches(0.18),
                 LEFT_W - Inches(0.35), content_h - Inches(0.25))

    # Right plain-English panel
    plain_english_panel(slide, plain_lines,
                        RIGHT_L, CONTENT_TOP,
                        RIGHT_W, content_h, accent=ac)

    # Constraint (if any)
    ct = content_bot + Inches(0.08)
    if constraint:
        constraint_box(slide, constraint,
                       Inches(0.3), ct, SLIDE_W - Inches(0.4))
        ct += constraint_h

    # Key Insight
    key_insight(slide, insight,
                Inches(0.3), ct, SLIDE_W - Inches(0.4), accent=ac)

    brand_bar(slide, num)
    return slide



# 
# SLIDES
# 

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_JET)

    # hero panel
    rect(slide, 0, 0, SLIDE_W, Inches(3.3), C_NAVY)
    left_bar(slide, C_ACCENT)

    badge(slide, "INTERNAL TECHNICAL REVIEW ",
    Inches(0.45), Inches(0.25), Inches(2.9), Inches(0.28),
    bg=C_ACCENT, fg=C_JET, size=8)

    txt(slide,
    "Predictive Project Management\nInformation System (PMIS)",
    Inches(0.45), Inches(0.65), Inches(11.5), Inches(1.8),
    size=38, bold=True, color=C_WHITE)

    txt(slide,
    "ML-Assisted Delay Forecasting | Alpha Prototype Review",
    Inches(0.45), Inches(2.48), Inches(11.5), Inches(0.5),
    size=16, italic=True, color=C_CYAN)

    rule(slide, Inches(0.45), Inches(3.28), SLIDE_W - Inches(0.6), color=C_ACCENT)

    # meta cards
    meta = [
    ("Project Lead", "Devesh Jha"),
    ("Phase", "Alpha Validation"),
    ("Focus", "ML Pipeline and Architecture"),
    ("Status", "Prototype: Logic Validated"),
    ]
    cw, ch = Inches(2.9), Inches(1.15)
    for i, (lbl, val) in enumerate(meta):
        cx = Inches(0.45) + i * Inches(3.1)
        cy = Inches(3.45)
        rect(slide, cx, cy, cw, ch, C_PANEL)
        rect(slide, cx, cy, cw, Inches(0.06), C_ACCENT)
        txt(slide, lbl.upper(),
        cx + Inches(0.15), cy + Inches(0.13),
        cw - Inches(0.3), Inches(0.28),
        size=8, bold=True, color=C_ACCENT)
        txt(slide, val,
        cx + Inches(0.15), cy + Inches(0.45),
        cw - Inches(0.3), Inches(0.6),
        size=14, bold=False, color=C_WHITE)

    # bottom ribbon
    rect(slide, 0, SLIDE_H - Inches(1.02),
         SLIDE_W, Inches(0.74), RGBColor(0x0A, 0x14, 0x1E))
    lines = [
        "WHAT IS THIS?  I have built an AI system that predicts when a construction project will be delayed - before it actually gets delayed.",
        "CURRENT STATUS  The core math and AI logic has been tested and validated. This review is to share findings and discuss the path forward.",
        "GOAL OF TODAY  Get alignment on current architecture trade-offs and discuss the Phase 2 stabilisation plan.",
    ]
    for i, line in enumerate(lines):
        txt(slide, line,
            Inches(0.5), SLIDE_H - Inches(1.0) + i * Inches(0.24),
            SLIDE_W - Inches(0.8), Inches(0.25),
            size=9, color=C_MUTED)

    brand_bar(slide, 1)
    return slide



def slide_02_problem(prs):
    return dual_slide(
    prs, 2, "", "Limitations of Deterministic CPM Tracking",
    "Why traditional project scheduling fails us in complex infrastructure projects",
    tech_bullets=[
    ("Static Duration Assumptions",
    "CPM assigns one fixed duration per task, assuming ideal conditions. A foundation task gets the same 30 days whether it is on a controlled urban site or a monsoon-season high-risk zone."),
    ("Omitted Environmental Variables",
    "Systems ignore LWE zone security friction, terrain difficulty, vendor reliability, and seasonal weather - factors that are the primary delay drivers in our portfolio."),
    ("Reactive, Not Predictive",
    "Delays are detected only AFTER they occur. A PM discovers in week 14 that the project is already 3 weeks behind, with cascading costs already materialised."),
    ("Subjective Buffer Application",
    "PMs add time buffers based on gut feel - one pads 10%, another 25% for the same scope. No standardised, mathematically defensible methodology exists."),
    ],
    plain_lines=[
    "Planning a road trip using only perfect traffic conditions - no rain, no detours. That is exactly what our current scheduling system does.",
    "We find out we are behind schedule AFTER we are already behind - like checking the fuel gauge after the car stalls.",
    "Different PMs add different safety buffers for the same work. Inconsistent and unscientific.",
    "The AI system solves all three of these problems systematically.",
    ],
    constraint="Current schedules are optimistic by design. No mathematical basis exists for the risk buffers in use today.",
    insight="Every major delay in our portfolio traces back to a predictable environmental variable - one that our legacy system was blind to.",
    accent=C_AMBER,
    )


def slide_03_solution(prs):
    return dual_slide(
    prs, 3, "", "Probabilistic Forecasting via Quantile Regression",
    "Replacing single-point estimates with ML-derived confidence intervals",
    tech_bullets=[
    ("Algorithm: CatBoost Quantile Regression",
    "Three CatBoost models are trained - one each for P10, P50, and P90 quantiles - replacing a single-point estimate with a statistically grounded confidence interval."),
    ("Output: P10 / P50 / P90 Confidence Bounds",
    "P10 = optimistic case, P50 = median expected completion, P90 = conservative planning buffer. Every task gets a range, not just one number."),
    ("Feature Integration: Environmental Friction",
    "The model auto-incorporates monsoon overlap (day-count), LWE zone friction, vendor tier scores, and terrain difficulty - combined at inference time."),
    ("Real Data Acquisition Path",
    "Current training uses synthetic data. Path to real data: (1) extract ERP task records, (2) map to 8-feature schema, (3) retrain. The pipeline is ready - only the data source changes."),
    ],
    plain_lines=[
    "Instead of '30 days', the system says: 'best case 25, most likely 30, but plan for 42 to be safe.'",
    "Like a weather forecast: '70% chance of rain' is more useful than 'it might rain'. Same logic applied to project timelines.",
    "The AI auto-checks: monsoon overlap? vendor reliability? high-risk zone? Each factor adjusts the forecast.",
    "Connecting to real ERP data is a well-defined next step - not a research problem.",
    ],
    constraint="Model trained on synthetic data, not live ERP history. Calibration to real data is the defined next milestone.",
    insight="P10/P50/P90 bounds give decision-makers a mathematically grounded basis for contingency budgeting - not guesswork.",
    accent=C_GREEN,
    )


def slide_04_architecture(prs):
    return dual_slide(
    prs, 4, "", "Current System Architecture (Alpha Build)",
    "Understanding the deliberate design trade-offs of the monolithic prototype",
    tech_bullets=[
    ("Technology Stack: Streamlit Monolith",
    "A single Streamlit Python application where the entire UI, data routing, and ML inference run in one process - enabling rapid iteration without a separate frontend layer."),
    ("Data Flow: Synchronous Execution",
    "Every user interaction re-runs the full Python script top-to-bottom. Feature engineering and CatBoost inference are sequential - correct, but no parallelism."),
    ("Memory Management: Model Caching",
    "@st.cache_resource holds the three CatBoost bundles (~4 MB each) in memory across sessions, avoiding a fresh disk read and ~2-3s deserialization on each inference call."),
    ("Critical Constraint: Single-Threaded Blocking",
    "Heavy operations (e.g., 500-node Monte Carlo simulation) block the single OS thread, freezing the UI for all users until completion. This is the primary architectural debt."),
    ],
    plain_lines=[
    "A one-person shop: serves one customer perfectly, but ten customers simultaneously all wait in a single queue.",
    "Everything runs in sequence: input → AI inference → render. Nothing runs in parallel.",
    "This was intentional - the right trade-off to validate ML logic fast. Not a mistake.",
    "Phase 2 directly addresses this by splitting the system into independent components.",
    ],
    constraint="Single-threaded, single-process. Concurrent users cause blocking and latency spikes. Not suitable for multi-user production.",
    insight="The monolithic architecture was the right choice for prototype validation - proving ML logic in weeks rather than months. Phase 2 dismantles it systematically.",
    accent=C_ACCENT,
    )


def slide_05_ml_pipeline(prs):
    return dual_slide(
    prs, 5, "", "ML Pipeline and Feature Engineering",
    "How the model is built, serialised, and deployed for inference",
    tech_bullets=[
    ("CatBoost + Bayesian Ridge Baseline",
    "Three CatBoost models (P10/P50/P90) use loss_function='Quantile:alpha=q'. A Bayesian Ridge baseline is trained in parallel for direct comparison. All four models are bundled into data/pmis_model.pkl via joblib."),
    ("Training / Inference Separation",
    "Training runs offline via scripts/model.py. The result is a single artifact (pmis_model.pkl) loaded at startup via @st.cache_resource. The live app contains zero training code - pure inference consumer."),
    ("Active Feature Set (8 Features)",
    "Project_Type, District, LWE_Flag, Task_Category, Land_Type, Vendor_Tier, Planned_Duration, Monsoon_Flag. Vendor_Tier rating: 1 = premium (L&T/Tata), 2 = mid-tier, 3 = local contractor."),
    ("Engineered but Inactive: TF-IDF NLP Layer",
    "TF-IDF vectorisation on Task_Name (max_features=100) is implemented in CELL 7b but commented out - disabled to maintain baseline stability and interpretability during validation."),
    ],
    plain_lines=[
    "The AI was trained offline and packaged into one file (pmis_model.pkl). The live app loads it and predicts immediately - no retraining at runtime.",
    "Four models trained: 3 AI quantile models (best/likely/worst) + 1 simpler statistical baseline to verify the AI is genuinely better.",
    "8 inputs: project type, district, task category, land type, vendor tier, planned duration, LWE risk, monsoon overlap - all real, available project record fields.",
    "A text-analysis layer (task descriptions) is built but switched off to keep outputs clean. It is the next planned upgrade.",
    ],
    constraint="TF-IDF NLP layer commented out (CELL 7b). MLflow not yet active - versioning is file-system based (timestamped folders).",
    insight="Replacing pmis_model.pkl is all that is needed to update the model - zero code changes required in the live application.",
    accent=C_ACCENT,
    )


def slide_06_fastapi(prs):
    return dual_slide(
    prs, 6, "", "Proposed Backend Decoupling (Phase 2 Architecture)",
    "Engineering the transition from monolith to microservice",
    tech_bullets=[
    ("The Core Problem: Tight Coupling",
    "ML inference, data transformation, and UI rendering all share one Python execution context. A slow inference call degrades UI; a crash takes down everything. No separation of concerns, no independent scaling."),
    ("Proposed Solution: FastAPI Microservice",
    "FastAPI under Uvicorn ASGI supports async/await, enabling concurrent request handling. The Streamlit UI becomes a thin client calling the API via HTTP - instead of direct Python function calls."),
    ("Data Contracts via Pydantic",
    "Pydantic input models enforce field types, required fields, and value ranges at every endpoint. Invalid payloads are rejected with 422 errors before reaching the inference engine."),
    ("Current Status",
    "FastAPI is not yet implemented - deliberately excluded in Phase 1 to prioritise ML logic validation. Full implementation is the milestone-one deliverable of the Phase 2 sprint."),
    ],
    plain_lines=[
    "The UI and AI engine are currently glued in one file. Phase 2 splits them into two systems - like separating a waiter (UI) from the kitchen (AI engine).",
    "Once behind FastAPI, multiple users can send requests simultaneously without blocking each other.",
    "Pydantic acts as a strict bouncer - garbage data is rejected before the AI engine ever sees it.",
    "This is the next 30-day sprint. Planned, designed, and ready to execute.",
    ],
    constraint="FastAPI NOT yet implemented. Input validation relies solely on Streamlit UI constraints - fragile and bypassable.",
    insight="API decoupling is the single highest-priority Phase 2 action - it unlocks concurrent access, independent scaling, and eliminates the input validation vulnerability.",
    accent=C_CYAN,
    )


def slide_07_production_readiness(prs):
    return dual_slide(
    prs, 7, "", "Codebase Maturation Status",
    "Engineering improvements applied to stabilise the Alpha prototype",
    tech_bullets=[
    ("Dependency Pinning: requirements.txt",
    "All 15 production dependencies are version-pinned (e.g., catboost==1.2.10, streamlit==1.54.0). Any environment reproduces the exact binary stack - eliminating 'works on my machine' failures."),
    ("Centralised Logging: app.log",
    "print() statements replaced with Python's logging module at the orchestrator level (app/main.py). Output goes to stdout and app.log on disk, with severity filtering (INFO/WARNING/ERROR) and timestamps."),
    ("Dynamic Inference in Mid-Project Tracker",
    "Hardcoded schedule placeholders replaced with live CatBoost P50 inference calls. The tracker now recalculates duration dynamically when the user changes environmental parameters."),
    ("Testing Status",
    "One CPM smoke test (tests/smoke_test.py) exists. No Pytest unit tests for ML inference or UI modules. No CI/CD pipeline. Manual regression testing is the current sole method."),
    ],
    plain_lines=[
    "All 15 library versions locked exactly - like pinning the exact tools used to build a bridge. Anyone who installs the project gets an identical toolset.",
    "Replaced scattered print() messages with a structured app.log file with timestamps and severity - like a proper flight recorder.",
    "Mid-project estimates now recalculate live using AI every time parameters change, instead of showing static numbers.",
    "Honest gap: one CPM test exists, but no automated AI or UI tests. Phase 2 introduces Pytest coverage.",
    ],
    constraint="No Pytest unit tests for ML inference or UI. One CPM smoke test exists. No CI/CD pipeline - manual regression testing only.",
    insight="Centralised logging + dependency pinning make the system diagnosable and reproducible - critical preconditions before any team handover.",
    accent=C_GREEN,
    )


def slide_08_validation(prs):
    return dual_slide(
    prs, 8, "", "Validation Methodology and Metrics",
    "How model quality was measured and what the results actually tell us",
    tech_bullets=[
    ("Metrics: MAE, R2, RMSE on 20% Hold-out",
    "Training pipeline (TEST_SIZE=0.2, RANDOM_STATE=42) logs MAE, R2, and RMSE for all four models on a held-out test split. Metrics are saved inside the model artifact dictionary."),
    ("Bayesian Ridge Comparison Baseline",
    "A BayesianRidge model is trained on the same split for direct comparison. If CatBoost is not clearly superior, the added complexity is not justified. Results are printed side-by-side."),
    ("SHAP Explainability on P50 Model",
    "SHAP TreeExplainer generates a global bar chart and beeswarm plot for the P50 model. Output correctly identifies Vendor_Tier and Monsoon_Flag as the primary delay drivers."),
    ("Incorrect Predictions: Feedback Loop",
    "Divergent predictions are logged in a feedback ledger, accumulated, and included in the next retraining cycle as additional ground-truth data - closing the learning loop without code changes."),
    ],
    plain_lines=[
    "Tested on data it had never seen (20% hold-out). Error numbers are real logged outputs, not estimates.",
    "A simpler statistical model (Bayesian Ridge) was trained alongside for comparison - if AI is not clearly better, added complexity is unjustified.",
    "SHAP explains WHY each prediction was made. It correctly flagged vendor quality and monsoon season as the biggest delay drivers.",
    "Wrong prediction? The actual outcome is logged and feeds the next retraining cycle. The system self-corrects over time.",
    ],
    constraint="Metrics are on synthetic data only. Results will shift on real ERP data - expected, and the retraining pipeline is already built to handle it.",
    insight="SHAP output aligning with domain knowledge (Vendor Tier and Monsoon as top drivers) confirms the model learned the right relationships.",
    accent=C_AMBER,
    )


def slide_09_scalability(prs):
    return dual_slide(
    prs, 9, "", "Proposed Scalability Strategy",
    "Infrastructure prerequisites for enterprise-grade multi-user deployment",
    tech_bullets=[
    ("Containerisation: Docker",
    "Frontend, FastAPI, PostgreSQL, and Redis containers orchestrated by Docker Compose with explicit base-image pinning (python:3.11-slim). Provides a direct, reproducible path to AWS ECS / Azure AKS."),
    ("State Persistence: PostgreSQL Migration",
    "Current system uses session_state and CSV files - a server restart wipes all data. Phase 3 migrates to PostgreSQL storing project records, prediction history, and session tokens."),
    ("Async Task Processing: Celery + Redis",
    "500-node Monte Carlo simulations currently block the main thread for 8-10 seconds. Celery offloads them to background workers; the UI returns immediately with a 'processing' status."),
    ("Current State: Zero Production Infrastructure",
    "No database, no container, no task queue. Runs entirely on local memory and CSV files. All of this is Phase 3 scope - after Phase 2 API decoupling is complete."),
    ],
    plain_lines=[
    "Docker packages the entire application into a container that runs identically on any server - laptop, AWS, Azure. No more 'works on my machine'.",
    "Server restart currently wipes all data. PostgreSQL stores project data persistently, surviving restarts and supporting multiple users.",
    "Heavy AI computation currently freezes the screen for everyone. Celery runs it in the background - submit the job, get a notification when ready.",
    "All Phase 3. Phase 2 decouples the API first, then Phase 3 builds the infrastructure beneath.",
    ],
    constraint="Runs entirely on local memory and CSV files. No database, no container, no task queue. Server restart wipes all session data.",
    insight="Docker de-risks the cloud deployment decision entirely - the same image runs identically on AWS, Azure, or on-premise.",
    accent=C_ACCENT,
    )


def slide_10_business_impact(prs):
    return dual_slide(
    prs, 10, "", "Strategic Value Formulation",
    "How probabilistic ML forecasting creates direct operational and financial impact",
    tech_bullets=[
    ("P90 Contingency Quantification",
    "Flat 10-15% portfolio-wide buffers replaced with mathematically derived worst-case timelines specific to each project's friction profile. Contingency reserves proportional to computed risk, not blanket padding."),
    ("Proactive Intervention: Leading Indicators",
    "Delay probability distributions are computed at the planning stage - before ground breaks - enabling pre-emptive intervention: alternate vendor selection, adjusted sequencing, early contingency activation."),
    ("Targeted Accountability via SHAP",
    "SHAP maps delay contributions: e.g., a 45-day slippage = 28 days (Vendor Tier 3) + 12 days (Monsoon) + 5 days (LWE friction). Gives management evidence-based justification for procurement decisions."),
    ("Business Value Dependency: Data Quality",
    "All value propositions are contingent on retraining with real ERP execution data. P10/P50/P90 outputs carry decision-grade validity only after real-data calibration."),
    ],
    plain_lines=[
    "'Add 10% to every budget' replaced with 'this project needs 18% due to vendor risk + monsoon timing, but that other one needs only 6%'. Precise, not blanket.",
    "Currently, delays are discovered in week 14. This system flags high-risk tasks at week 1 planning - before anything is committed.",
    "AI predicts 45-day delay: SHAP shows '28 days - Tier 3 vendor, 12 days - monsoon, 5 days - LWE friction'. Concrete, data-backed reason to upgrade the vendor.",
    "Full potential requires retraining on real historical project data - the Phase 3 milestone we are building towards.",
    ],
    constraint="Full business value contingent on ERP data integration and real-data model retraining. Current outputs are indicative, not decision-grade.",
    insight="SHAP attribution transforms this from a black-box AI into a transparent decision tool - giving leadership data-backed justification for procurement and scheduling.",
    accent=C_GREEN,
    )


def slide_11_risks(prs):
    return dual_slide(
    prs, 11, "", "System Constraints and Failure Scenarios",
    "Transparent assessment of current architectural risks and known failure modes",
    tech_bullets=[
    ("Concurrency Failure: Single-Threaded Blocking",
    "Monte Carlo and CPM simulations block the single OS thread for 8-12 seconds. All other active users see a frozen UI during this window. Under 3+ concurrent heavy users, cascading timeouts occur."),
    ("Input Vulnerability: No Backend Data Contract",
    "Validation relies solely on Streamlit UI widgets. Bypassing the UI - via session state mutation or scripting - sends malformed data directly to the CatBoost engine, causing an unhandled exception and server crash."),
    ("Feature Completeness: NLP Layer Disabled",
    "TF-IDF vectorisation on task descriptions is implemented but disabled at inference time. Semantic signals in unstructured text are currently ignored - a documented next improvement."),
    ("Production Suitability",
    "Explicitly scoped as a single-user local validation tool. Concurrent access, input validation robustness, and persistent state are all Phase 2/3 prerequisites before any network exposure."),
    ],
    plain_lines=[
    "A toll booth with one lane. One car is fine; ten cars and nine wait - some give up. That is our system with multiple simultaneous users.",
    "The UI form is the only door. If someone finds a back door, garbage data crashes the AI engine. Phase 2 (FastAPI + Pydantic) builds a proper checkpoint.",
    "Text analysis feature built but switched off - deliberately, to keep outputs stable and explainable during this validation phase.",
    "This is a validated prototype, not a deployed product. Every constraint here has a planned resolution in the roadmap.",
    ],
    constraint="Multi-user concurrent access causes blocking and potential crashes. UI bypass can inject malformed data directly into the inference engine.",
    insight="Transparent disclosure of these failure modes demonstrates engineering maturity. Every constraint has a concrete, planned resolution in the Phase 2/3 roadmap.",
    accent=C_RED,
    )


def slide_12_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_JET)
    slide_header(slide, 12, "", "Technical Execution Roadmap",
    "Phased transition from Alpha prototype to production-ready connected system", C_ACCENT)

    phases = [
    {
    "num": "01",
    "label": "Logic Validation",
    "tag": "COMPLETE",
    "color": C_GREEN,
    "items": [
    ("Prototype Delivered", "Functional Streamlit monolith with full UI flow and interactive outputs."),
    ("ML Pipeline Validated", "CatBoost Quantile Regression for P10/P50/P90 inference confirmed stable on synthetic dataset."),
    ("SHAP Integrated", "Feature attribution working, correctly identifying Vendor Tier and Monsoon as primary friction drivers."),
    ("Codebase Stabilised", "Logging, dependency pinning, and dynamic inference substitution completed."),
    ],
    },
    {
    "num": "02",
    "label": "Architecture Decoupling",
    "tag": "NEXT PHASE",
    "color": C_AMBER,
    "items": [
    ("FastAPI Microservice", "Stand up FastAPI + Uvicorn async backend absorbing all ML inference and data transformation calls."),
    ("Pydantic Data Contracts", "Strict HTTP schema validation - reject malformed payloads before reaching inference engine."),
    ("Docker Containerisation", "Package frontend and API containers for environment-agnostic deployment."),
    ("Pytest Integration", "Establish automated unit test coverage for all inference endpoints and data transforms."),
    ],
    },
    {
    "num": "03",
    "label": "Async Scaling & ERP Data",
    "tag": "EVALUATING",
    "color": C_MUTED,
    "items": [
    ("Celery + Redis Queues", "Offload heavy Monte Carlo / CPM computations to async worker processes."),
    ("PostgreSQL Migration", "Replace in-memory session state and CSVs with a persistent relational schema."),
    ("ERP Data Ingestion", "Map, normalise, and ingest real historical project execution data to trigger true-variance model retrain."),
    ("MLflow Tracking", "Formal experiment tracking, model versioning, and data drift monitoring."),
    ],
    },
    ]

    col_w = Inches(3.9)
    gap = Inches(0.2)
    x0 = Inches(0.3)
    top = Inches(1.62)
    ch = Inches(5.35)
    colors_map = {0: C_GREEN, 1: C_AMBER, 2: C_MUTED}

    for i, ph in enumerate(phases):
        cx = x0 + i * (col_w + gap)
        ac = colors_map[i]
        rect(slide, cx, top, col_w, ch, C_PANEL)
        rect(slide, cx, top, col_w, Inches(0.08), ac)

        txt(slide, f"PHASE {ph['num']}",
        cx + Inches(0.15), top + Inches(0.14),
        col_w - Inches(0.3), Inches(0.28),
        size=9, bold=True, color=ac)

        badge(slide, ph["tag"],
        cx + Inches(0.15), top + Inches(0.45),
        col_w - Inches(0.3), Inches(0.28),
        bg=ac, fg=C_JET, size=8)

        txt(slide, ph["label"],
        cx + Inches(0.15), top + Inches(0.82),
        col_w - Inches(0.3), Inches(0.45),
        size=15, bold=True, color=C_WHITE)

        rule(slide, cx + Inches(0.15), top + Inches(1.3),
        col_w - Inches(0.3), color=ac)

        for j, (bl, bt) in enumerate(ph["items"]):
            yy = top + Inches(1.42) + j * Inches(0.95)
            tb = slide.shapes.add_textbox(
            cx + Inches(0.15), yy, col_w - Inches(0.3), Inches(0.92))
            tf = tb.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            r1 = p1.add_run()
            r1.text = f"{bl}"
            r1.font.bold = True
            r1.font.size = Pt(10.5)
            r1.font.color.rgb = ac
            r1.font.name = "Calibri"
            p2 = tf.add_paragraph()
            p2.space_before = Pt(2)
            r2 = p2.add_run()
            r2.text = bt
            r2.font.size = Pt(9.5)
            r2.font.color.rgb = C_WHITE
            r2.font.name = "Calibri"

    key_insight(slide, "Phase 2 FastAPI decoupling is the critical-path milestone. Every Phase 3 capability - async scaling, ERP data, persistent state - depends on the API layer being in place first.",
                Inches(0.3), SLIDE_H - Inches(0.27) - Inches(0.65),
                SLIDE_W - Inches(0.4), accent=C_ACCENT)

    brand_bar(slide, 12)
    return slide


def slide_13_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_JET)
    slide_header(slide, 13, "", "Prototype Capabilities Demonstration",
    "Live navigation of the PMIS Alpha build - two-scenario walkthrough", C_ACCENT)

    # hero box
    rect(slide, Inches(3.2), Inches(1.65),
    Inches(6.8), Inches(1.7), C_PANEL)
    txt(slide, "LIVE ALPHA BUILD DEMO",
    Inches(3.2), Inches(1.85), Inches(6.8), Inches(1.2),
    size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    txt(slide, "Run: streamlit run pmisapp.py",
    Inches(3.2), Inches(3.05), Inches(6.8), Inches(0.35),
    size=11, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    steps = [
    ("Step 1\nBaseline",
    "Set parameters: Tier 1 Vendor · Low LWE Risk District · April start (pre-monsoon).\n"
    "→ Observe P50 duration as the ideal-condition baseline reference."),
    ("Step 2\nFriction",
    "Modify: switch to Tier 3 Vendor · High LWE Risk Zone · June start (deep monsoon).\n"
    "→ Observe P90 bound stretch dramatically - each parameter multiplies delay risk."),
    ("Step 3\nSHAP",
    "Open SHAP explainability panel for the friction scenario.\n"
    "→ Review per-feature attribution: which factor contributed most days of delay."),
    ("Step 4\nCPM",
    "Navigate to the Critical Path view.\n"
    "→ Observe how the delay distributes across the task network and shifts the project end date."),
    ]

    sw = Inches(2.9)
    for i, (lbl, desc) in enumerate(steps):
        cx = Inches(0.3) + i * Inches(3.25)
        cy = Inches(3.5)
        sh = Inches(2.75)
        rect(slide, cx, cy, sw, sh, C_PANEL)
        rect(slide, cx, cy, sw, Inches(0.06), C_ACCENT)
        badge(slide, lbl, cx, cy - Inches(0.35), sw, Inches(0.3),
        bg=C_ACCENT, fg=C_JET, size=8)
        txt(slide, desc,
            cx + Inches(0.12), cy + Inches(0.14),
            sw - Inches(0.22), sh - Inches(0.22),
            size=10.5, color=C_WHITE)

    key_insight(slide, "The two-scenario walkthrough makes the AI's sensitivity to environmental friction immediately tangible - even for non-technical stakeholders watching the numbers change in real time.",
                Inches(0.3), SLIDE_H - Inches(0.27) - Inches(0.65),
                SLIDE_W - Inches(0.4))
    brand_bar(slide, 13)
    return slide


def slide_14_qa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_JET)
    slide_header(slide, 14, "", "Appendix - Technical Review & Q&A",
    "Anticipated questions and prepared responses for the review panel", C_ACCENT)

    qas = [
    ("Is this production ready?",
    "No - deliberately. This Alpha build prioritises ML logic correctness. FastAPI decoupling and Pydantic validation must be completed before exposing to network traffic or concurrent users."),
    ("Why Streamlit and not a proper web framework?",
    "Streamlit converts Python data science code into interactive UI without frontend engineering - enabling rapid iteration on ML outputs. It is a temporary presentation layer, not the final architecture."),
    ("What happens if multiple users access it simultaneously?",
    "The server thread blocks and serves sequentially. Under 3+ concurrent heavy users, response times spike and later requests timeout or receive 503 errors. This is the primary Phase 2 target."),
    ("If trained on fake data, why trust it at all?",
    "The specific weights are uncalibrated, but the pipeline - feature ingestion, quantile inference, SHAP, CPM integration - is architecturally sound. Switching to real ERP data requires only a data source change and retrain."),
    ("Are you tracking model performance over time?",
    "Not yet. MLflow tracking and drift alerting are Phase 3 scope. Currently, versioning uses timestamped file-system artifact folders - functional but not enterprise-grade."),
    ("What do you need from leadership today?",
    "1) Approve the Phase 2 FastAPI sprint (30-day scope).  2) Create a channel for ERP data access discussions.  3) Acknowledge the architectural limitations - this is not production-ready as-is."),
    ]

    cw = (SLIDE_W - Inches(0.75)) / 2
    for i, (q, a) in enumerate(qas):
        col = i % 2
        row = i // 2
        cx = Inches(0.3) + col * (cw + Inches(0.15))
        cy = Inches(1.65) + row * Inches(1.62)
        card_h = Inches(1.52)
        rect(slide, cx, cy, cw, card_h, C_PANEL)
        rect(slide, cx, cy, Inches(0.06), card_h, C_ACCENT)
        txt(slide, f"Q: {q}",
        cx + Inches(0.15), cy + Inches(0.08),
        cw - Inches(0.2), Inches(0.38),
        size=10.5, bold=True, color=C_ACCENT)
        txt(slide, a,
            cx + Inches(0.15), cy + Inches(0.45),
            cw - Inches(0.2), Inches(1.0),
            size=10, color=C_WHITE)

    key_insight(slide, "The strongest answer to any tough question today: 'I know the limitation. It was built this way deliberately to validate the ML logic fast. The roadmap addresses it directly.'",
                Inches(0.3), SLIDE_H - Inches(0.27) - Inches(0.65),
                SLIDE_W - Inches(0.4))
    brand_bar(slide, 14)
    return slide


# 
# MAIN
# 

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
    ("Title", slide_01_title),
    ("Problem", slide_02_problem),
    ("Solution", slide_03_solution),
    ("Architecture", slide_04_architecture),
    ("ML Pipeline", slide_05_ml_pipeline),
    ("FastAPI", slide_06_fastapi),
    ("Production Readiness", slide_07_production_readiness),
    ("Validation", slide_08_validation),
    ("Scalability", slide_09_scalability),
    ("Business Impact", slide_10_business_impact),
    ("Risks", slide_11_risks),
    ("Roadmap", slide_12_roadmap),
    ("Demo", slide_13_demo),
    ("Q&A", slide_14_qa),
    ]

    for i, (name, fn) in enumerate(builders, 1):
        fn(prs)
        print(f"Slide {i:02d} - {name}")

        out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
        "PMIS_Alpha_Review.pptx")
        prs.save(out)
        print(f"\n Saved → {out}")
        print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
